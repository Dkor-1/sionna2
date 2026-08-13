# -*- coding: utf-8 -*-
"""
make_0811_v20.py — **v19 에서 v20 을 짓는다.** 두 가지만 한다.

  ① 발표 대본(notes)을 전부 갈아끼운다 — 사실 오류 다섯 건 정정 + 짧고 쉬운 문장으로.
  ② **Future work 바로 앞에 슬라이드 두 장**을 새로 끼운다 —
     [9] 광선 예산(1.78 억 발 ↔ 40 억 발, 시드 2 개씩) · [10] 40 억 발 ↔ 우리 PO 나란히.

사용자(2026-08-11)
> "v20 만들 때 방금 만든 raybudget_40m 실어서. 대본은 팀미팅 직전에 시뮬레이션이 끝난
>  사항인데, 쏘는 ray 의 개수를 최대치로 늘렸더니 백그라운드가 훨씬 깨끗해졌다라고 말하고 싶어.
>  future work 로 가기 직전 슬라이드에 넣어서 구성해줘."
> "발표자 노트는 아주 간략히 — ray 개수를 늘려서 시뮬레이션 돌리면 오래 걸려서
>  방금 전 확인 가능했다 정도만."

⛔ v19 는 **안 건드린다.** 새 파일 v20 을 만든다.
⛔ 기존 아홉 장의 **그림·본문·레이아웃은 한 글자도 안 바꾼다.** (쪽번호 /9 → /11 만 예외)

■ 새 슬라이드를 어떻게 만드나
슬라이드 8("TWICE")을 **XML 통째로 복제**해 틀로 쓴다. 새로 그리지 않는 이유는 이 덱의
활자·색·여백이 빌더가 아니라 파일 안에 있어서, 손으로 다시 만들면 미세하게 어긋나기
때문이다. 복제 뒤 그림만 갈아 끼우고 세 군데 글자(태그·제목·결론띠)를 바꾼다.
그림은 틀 안에 **비율 유지로** 넣고 가로·세로 모두 가운데 정렬한다 —
가로로 넓은 것은 폭이, 덜 넓은 것은 높이가 먼저 걸린다.

■ 정정한 대본 다섯 (전부 사실 오류)
 1. [3] "hard-coded **randomized** values" — 난수가 아니다. 원장 `rpm_per_rotor` 는
    3800 × [+0.220 %, −0.220 %, −0.121 %, +0.121 %] 로 손으로 정한 결정론적 패턴이다.
 2. [4] "under a **plane-wave** illumination assumption" — 덱의 거리별 그림은 구면파다
    (`outputs/deck0811_range_figs.json` → `ours_illumination = "spherical wave, computed per range"`).
 3. [5] "127 Hz ≈ **propeller rotation speed**" — 2 배 틀렸다. 126.67 Hz 는 **블레이드 플래시율**
    (날개 2 장 × 3,800 rpm ÷ 60)이고 회전수는 63.3 Hz 다.
 4. [5] "밴드 430–1229 Hz 인데 피크가 127 Hz" — 그대로 들으면 자기모순이다. 대역 전력을
    **시간축으로 다시 FFT** 한 것이라는 두 단계를 명시했다.
 5. [6] "patterns remained highly consistent" — 우리 커널에서만 참이다.
 ⭐추가(사용자 지시): [3] 은 3,800 의 계산 과정을 **설명하지 않는다**. «계산값이지
    실측이 아니다» 까지만 말한다. [4] 는 "isotropically" 표현을 유지한다.
    [5] 는 CW 를 쓴 이유(초기 실험·다루기 쉬움·목적이 마이크로도플러 자체)를 넣는다.

    /workspace/.venvs/py312/bin/python teammeeting_0811/make_0811_v20.py
"""
from __future__ import annotations

import copy
import os
import shutil

from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches

TM = "/workspace/team_meeting"
SIONNA2 = "/workspace/sionna"
SRC = f"{TM}/teammeeting_0811_v19.pptx"
DST = f"{TM}/teammeeting_0811_v20.pptx"

TEMPLATE_IDX = 7          # 0-기준. 슬라이드 8 "TWICE" — 제목 + 큰 그림 + 결론띠
INSERT_AT = 8             # 0-기준. Future work(현 index 8) 바로 **앞**부터 차례로 끼운다

# 틀 슬라이드의 글자 세 자리 — 이것을 열쇠로 갈아 끼운다
K_TAG = "TWICE"
K_TITLE = "Forty metres, repeated runs"
K_BAR = "Only the random start differs between these runs."

# 끼워 넣을 슬라이드 — INSERT_AT 부터 이 순서대로 들어간다
NEW_SLIDES = [
    dict(tag="RAY BUDGET",
         title="More rays, cleaner background",
         bar="Pushing the ray budget to its maximum cleaned up the background.",
         fig=f"{SIONNA2}/outputs/figures/raybudget_40m_deck.png"),
    dict(tag="SIDE BY SIDE",
         title="The two engines at forty metres",
         bar="Same distance, same drone. Ours repeats exactly, the Path Solver does not.",
         fig=f"{SIONNA2}/outputs/figures/raybudget_40m_vs_ours.png"),
]

# 최종 11 장 기준 대본
NOTES = {
    2: "Today I show micro-Doppler maps of a hovering drone, computed two ways: "
       "with Sionna's Path Solver, and with our own SBR plus physical optics kernel.",

    3: "The target is a DJI Matrice 4E. Each rotor spins at a fixed speed near 3,800 RPM. "
       "That is a calculated value, not something I measured. The four rotors differ slightly "
       "from each other, and that spread comes from Ji-hyuk's prior work. These are hand-set "
       "constant values, not taken from a real flight. Better rotor modelling is future work.",

    4: "Sionna's Path Solver launches rays isotropically from the transmitter, that is, "
       "evenly in every direction. When a ray hits a surface it picks up that material's "
       "reflection coefficient, and the result comes from the few paths that actually "
       "come back.\n\n"
       "Our kernel does something different. We shoot a grid of parallel rays at the target, "
       "take the first point each ray hits, and sum the physical optics scattering integral "
       "over those points. For parts we treat as transmissive, the ray passes through, "
       "and we add the echo from the metal inside, scaled by how much gets through.",

    5: "The centre frequency is 3.5 GHz, a standard 5G band. For the waveform I used a simple "
       "continuous wave. This is an early experiment, and CW is the easiest thing to work "
       "with. What I want to see first is the micro-Doppler itself, so I did not need a more "
       "complicated waveform yet.\n\n"
       "For the STFT: window 70 samples, about 3.6 milliseconds; hop 2 samples; "
       "FFT size 560.\n\n"
       "The bottom figure is made in two steps. First I take the energy in the 430 to "
       "1229 Hz band, which is above the body and belongs to the blade tips. Then I ask "
       "how that energy rises and falls over time, and take its spectrum.\n\n"
       "The peak sits at about 127 Hz. That is the blade flash rate: two blades passing per "
       "revolution, so twice the rotation rate. Its harmonics show up clearly too.",

    6: "The previous slide was 3 metres only. Here I ran the same thing at 3, 15 and "
       "40 metres.\n\n"
       "Our kernel gives nearly the same map at every range. The Path Solver does not. "
       "At 40 metres its map changes, and I will come back to why.",

    7: "This is the band energy for each range. One thing stands out in the harmonics. "
       "With our PO approach, the peak height gradually decreases as we move to the higher "
       "harmonics. With the Path Solver, the second harmonic often comes out almost as high "
       "as the first one, and sometimes higher.",

    8: "At 40 metres I ran the Path Solver again, changing only the random seed. "
       "Even with the same drone, the same pose and the same orientation, the seed changes "
       "the result. One seed puts the strongest line on the second harmonic instead of "
       "the first one.\n\n"
       "Our PO approach has no randomness. Same pose, same orientation, "
       "same answer every time.",

    # ⭐새 슬라이드 둘 — 사용자 지시대로 아주 짧게
    9: "One last thing. I raised the ray count to the maximum the solver allows. "
       "It takes a long time to run, so I only got the result just before this meeting. "
       "The background is much cleaner.\n\n"
       "[if asked] The seed still decides which harmonic comes out strongest.",

    10: "And here they are next to each other at 40 metres. Two Path Solver runs with the "
        "full ray budget, and our PO result.\n\n"
        "[if asked] The blade lines land in the same place. What changes between the two "
        "Path Solver runs is only the random seed.",

    11: "For my future work, I want to make the simulation more realistic. "
        "I have three directions.\n\n"
        "First, the rotor speeds. Right now each rotor turns at a hard-coded value, and the "
        "spread between them comes from a flight simulator, not from a real drone. I want to "
        "use real flight log data instead, and also introduce temporal variation, the way a "
        "real drone adjusts its rotors while it stays in place.\n\n"
        "Second, noise. There is no receiver noise in this simulation. That is why nothing "
        "became harder when the drone moved further away. The echo does get weaker with "
        "distance, but each map is scaled to its own brightest point, so we cannot see that. "
        "And with no noise, nothing is ever buried. In the real world, noise is what decides "
        "how far we can detect, so I want to add it.\n\n"
        "Third, the geometry and the environment. Everything today uses one station that "
        "sends and receives at the same place, in empty space. Our real experiment is a "
        "passive bistatic setup, where the signal comes from someone else's base station, "
        "and it happens in a room with walls and a floor. I want to move to that geometry "
        "and add clutter modelling.\n\n"
        "All three will make the results worse, and that is the point. "
        "They bring the simulation closer to a real measurement.",
}


def set_text(tf, s: str) -> None:
    """첫 run 의 서식을 살린 채 글자만 바꾼다. `.text =` 는 서식을 날린다."""
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)
    p = tf.paragraphs[0]
    runs = list(p.runs)
    if not runs:
        p.text = s
        return
    runs[0].text = s
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def clone_slide(prs, index: int):
    """슬라이드를 XML 통째로 복제해 맨 뒤에 붙인다. **그림(p:pic)은 빼고** 복제한다.

    p:pic 은 r:embed 로 이미지 파트를 가리키는데 그 관계는 슬라이드마다 따로다.
    어차피 그림을 갈아 끼울 것이므로 복제 대상에서 빼고 나중에 add_picture 로 새로 넣는다.
    """
    src = prs.slides[index]
    dst = prs.slides.add_slide(src.slide_layout)
    for shp in list(dst.shapes):                       # 레이아웃이 깔아 놓은 것 비우기
        shp._element.getparent().remove(shp._element)
    tree = dst.shapes._spTree
    for shp in src.shapes:
        if shp.shape_type == 13 or shp._element.tag.endswith("}pic"):
            continue                                   # 그림은 건너뛴다
        tree.append(copy.deepcopy(shp._element))
    return dst


def move_slide(prs, frm: int, to: int) -> None:
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[frm])
    lst.insert(to, ids[frm])


def main() -> None:
    need = [SRC] + [s["fig"] for s in NEW_SLIDES]
    for f in need:
        if not os.path.exists(f):
            raise SystemExit(f"❌ {f} 가 없다")
    shutil.copy2(SRC, DST)
    prs = Presentation(DST)
    n0 = len(prs.slides)
    print(f"═══ v19({n0} 장) → v20 · 새 슬라이드 {len(NEW_SLIDES)} 장\n")

    tmpl = prs.slides[TEMPLATE_IDX]
    pic_box = next(((s.left, s.top, s.width, s.height) for s in tmpl.shapes
                    if s.shape_type == 13), None)
    if pic_box is None:
        raise SystemExit("❌ 틀 슬라이드에 그림이 없다")
    L, T, W, H = pic_box

    # ── ① 새 슬라이드들 ─────────────────────────────────────────────────────
    for k, spec in enumerate(NEW_SLIDES):
        new = clone_slide(prs, TEMPLATE_IDX)

        # 그림 — 틀 **안에** 비율 유지로 넣고 가로·세로 모두 가운데.
        #   폭이 걸리면 폭에, 높이가 걸리면 높이에 맞춘다.
        iw, ih = Image.open(spec["fig"]).size
        ar = iw / ih
        if W / H > ar:                                 # 높이가 먼저 걸린다
            h, w = H, int(round(H * ar))
        else:                                          # 폭이 먼저 걸린다
            w, h = W, int(round(W / ar))
        left, top = L + (W - w) // 2, T + (H - h) // 2
        new.shapes.add_picture(spec["fig"], left, top, width=w, height=h)

        swap = {K_TAG: spec["tag"], K_TITLE: spec["title"], K_BAR: spec["bar"]}
        hit = 0
        for shp in new.shapes:
            if not shp.has_text_frame:
                continue
            cur = shp.text_frame.text.strip()
            if cur in swap:
                set_text(shp.text_frame, swap[cur])
                hit += 1
        if hit != len(swap):
            raise SystemExit(f"❌ 글자 교체 {hit}/{len(swap)} — 틀이 바뀌었다")

        move_slide(prs, len(prs.slides) - 1, INSERT_AT + k)
        print(f"  [{INSERT_AT + k + 1:>2}] {spec['tag']:<12} {spec['title']}")
        print(f"       {os.path.basename(spec['fig'])}  {iw}x{ih} ({ar:.2f}:1) → "
              f"W{Emu(w).inches:.2f} H{Emu(h).inches:.2f} in "
              f"(틀 W{Emu(W).inches:.2f} H{Emu(H).inches:.2f})")

    # ── ② 쪽번호 총장수 "/9" → "/10" ────────────────────────────────────────
    n1 = len(prs.slides)
    fixed = 0
    for slide in prs.slides:
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            for p in shp.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip() == f"/{n0}":
                        r.text = f"/{n1}"
                        fixed += 1
    print(f"  쪽수  «/{n0}» → «/{n1}»  {fixed} 곳")

    # ── ③ 대본 ─────────────────────────────────────────────────────────────
    print()
    for i, slide in enumerate(prs.slides, start=1):
        new_note = NOTES.get(i)
        if not new_note:
            print(f"  [{i:>2}] (대본 없음)")
            continue
        old = (slide.notes_slide.notes_text_frame.text or "").strip()
        slide.notes_slide.notes_text_frame.text = new_note
        print(f"  [{i:>2}] {len(old):>4d} → {len(new_note):>4d} 자   "
              f"{new_note.splitlines()[0][:56]}…")

    prs.save(DST)

    # ── 검산 ───────────────────────────────────────────────────────────────
    a, b = Presentation(SRC), Presentation(DST)
    inserted = set(range(INSERT_AT, INSERT_AT + len(NEW_SLIDES)))
    order = [i for i in range(n1) if i not in inserted]     # 새 장을 뺀 나머지

    def body(slide):
        """본문 글자만. **쪽번호 자리표시자는 뺀다** — «/9 → /10» 은 의도한 변경이다."""
        out = []
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if getattr(sh, "is_placeholder", False) \
                    and str(sh.placeholder_format.type).startswith("SLIDE_NUMBER"):
                continue
            out.append(sh.text_frame.text)
        return out

    diffs = [(j + 1, k + 1) for j, k in enumerate(order)
             if body(a.slides[j]) != body(b.slides[k])]
    same_body = not diffs
    if diffs:
        print(f"\n  ⚠ 본문이 다른 장(v19→v20): {diffs}")
        for j, k in diffs[:2]:
            for x, y in zip(body(a.slides[j - 1]), body(b.slides[k - 1])):
                if x != y:
                    print(f"     v19[{j}] {x[:70]!r}\n     v20[{k}] {y[:70]!r}")
    notes_ok = all(b.slides[i - 1].notes_slide.notes_text_frame.text.strip() == v.strip()
                   for i, v in NOTES.items())
    n_pic = [sum(1 for sh in b.slides[i].shapes if sh.shape_type == 13)
             for i in sorted(inserted)]

    print(f"\n  장수 {n0} → {n1} · 새 장 그림 {n_pic}"
          f" · 기존 장 본문 불변 {'✅' if same_body else '❌'}"
          f" · 대본 되읽기 {'✅' if notes_ok else '❌'}")
    print(f"\n✅ {DST}  ({os.path.getsize(DST):,} B)")
    if not (same_body and notes_ok and n_pic == [1] * len(NEW_SLIDES)):
        raise SystemExit("❌ 검산 실패 — 이 v20 을 쓰지 마라")


if __name__ == "__main__":
    main()
