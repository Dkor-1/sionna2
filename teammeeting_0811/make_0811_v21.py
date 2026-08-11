# -*- coding: utf-8 -*-
"""
make_0811_v21.py — **v20 에서 대본 세 군데만 고친다.** 슬라이드·그림·본문은 손대지 않는다.

사용자(2026-08-11)
> "현재 팀미팅 자료에서 **새로운 정보를 넣지 않는 선에서 대본만 수정**해서 v21 제작도 들어가줘"

고친 근거: 16 에이전트 적대적 검증(워크플로 w5wc803o7). 각 지적을 다시 반증 시도해
살아남은 것만 반영했다. 원장·Sionna 설치본 소스를 직접 읽어 확인된 것들이다.

■ ⛔ [6] 가장 심각 — 두 슬라이드 뒤 자기 자신과 충돌한다
  전: "The Path Solver does not. **At 40 metres its map changes**, and I will come back to why."
  이 문장은 «거리 때문»으로 들린다. 원장은 그렇게 말하지 않는다:
   · 40 m 시드 3 판의 박자 253.27 / 126.51 / 128.13 Hz — **셋 중 둘은 예측(126.67)을 맞춘다**
   · 그런데 슬라이드 6 의 40 m 칸이 하필 그 이상판이다
     (`build_deck0811_range_figs.py:63` → `SZ["S1/E"]`)
   · 3 m·15 m 는 각각 **한 판뿐**이고 seed=1 하드코딩
     (`report07_three_engine_ranges.py:82`) → 거리와 시드가 원장에서 **분리되지 않는다**
   · 맵 상관: PathSolver 3↔15 = **0.688**(모든 쌍 중 최저) · 40 m 시드쌍 = 0.812~0.870
     ⇒ **40 m 시드끼리가 3 m↔15 m 보다 더 닮았다.** 거리 효과 < 시드 효과
   · 40 억 발에서도 seed1 253.03 / seed2 126.22 → «광선이 모자랐다» 해명도 막힌다
   · 그리고 슬라이드 8 이 "changing only the random seed" 라고 말한다 — **정면 충돌**
  ✅ 반대로 «우리 커널은 거리마다 거의 같다» 는 튼튼하다 (상관 0.983 / 0.999 / 0.983)

■ ⚠ [4] 주사위가 어디서 던져지는지 안 말한다
  전: "When a ray hits a surface it picks up that material's reflection coefficient"
  실제로는 히트마다 **거울반사냐 확산이냐를 추첨**한다
  (`radio_material.py:406` `_sample_event_type`, `:811-813` 분기, `:764-766` prs·prd).
  발사 방향은 Fibonacci lattice 라 **결정론적**이고(`ray_tracing.py:16`), 시드가 들어가는 곳은
  **오직 이 추첨과 확산 방향**뿐이다(`sb_candidate_generator.py:110, :364-366`).
  ⭐그리고 `DRONE_GROUP_MAT` 에서 **prop = prop_plastic (S = 0.20)** 이라 주사위가
  **프로펠러 위에서 던져진다** — 변조를 만드는 바로 그 부위다. 이것이 시드가 답을 바꾸는 이유다.
  지금 [4] 대본에 확률어가 **한 개도 없어서**, [8] 의 "random seed" 가 갑자기 튀어나온다.

■ ⚠ [5] "belongs to the blade tips" 가 과장이다 — 그리고 ⭐**대역은 바꾸지 않기로 했다**
  창 70 샘플 → 빈 281.4 Hz, Hann 주엽 반폭 **563 Hz**. 대역 하단 430 Hz 는 **그 안쪽**이라
  0-도플러 동체선의 창 스커트가 샌다(422 Hz 에서 동체 첨두 대비 −15.3 dB, Hann 누설
  이론값과 1 dB 이내 일치).

  ⭐**그래서 하단을 널(563 Hz)로 올려 보려 했고, 재보니 안 된다.** 실측:

    셀                 430~1229 (현)      563~1229 (A)
    Sionna 3 m         126.26 Hz  ✅      **254.60 Hz  ❌ 2 배선으로 뒤집힘**
    Sionna 15 m        126.74             126.72
    Sionna 40 m s1     253.27             253.42
    Sionna 40 m s2     126.51             126.49
    우리 3/15/40 m     126.51/126.11/126.05   전부 동일

  Sionna 3 m 은 1× 가 2× 를 **+1.27 dB** 로 겨우 이기고 있어(A 에서 −0.35 dB) 대역을 조금만
  건드려도 넘어간다. A 로 가면 덱의 헤드라인 칸 하나가 망가진다. ⇒ **현행 대역 유지.**

  ■ 그리고 «0 근처가 움직인다» 는 관찰이 맞다 (사용자 지적)
  주파수 빈마다 정지분(시간평균)과 뛰는분(시간분산)을 갈라 재면:
    · 0~563 Hz 가 맵 전체 변조의 **98 %** 를 갖는다. 지금 대역은 그 중 **1.7 %** 만 본다.
    · 다만 그 맥동은 «동체 + 블레이드» 의 **교차항**과 **날개 안쪽**(도플러 ∝ 반경, 563 Hz 는
      f_tip 의 46 % 이므로 날개 안쪽 절반이 그 아래)이 섞인 것이다.
    · 누설 구간(430~563)은 **정지분 36 % · 뛰는분 12 %** 로 DC 가 무겁다.
  ⇒ 대역을 «약해서 이걸 본다» 가 아니라 **«동체를 일부러 피한다»** 로 설명해야 오해가 없다.
  + "blade tips" 도 부정확 — 0.35~1.0 f_tip 은 날개 **바깥 65 %** 다.

■ 손대지 않은 것
  슬라이드 11 장·그림 11 개·본문 글자·쪽번호 — 전부 v20 그대로. 대본 세 개만 다르다.

    /home/yunjung/.venvs/py312/bin/python teammeeting_0811/make_0811_v21.py
"""
from __future__ import annotations

import hashlib
import os
import shutil
import zipfile

from pptx import Presentation

TM = "/home/yunjung/workspace/team_meeting"
SRC = f"{TM}/teammeeting_0811_v20.pptx"
DST = f"{TM}/teammeeting_0811_v21.pptx"

# 갈아 끼울 대본만. 나머지 슬라이드는 v20 것을 그대로 둔다.
NOTES = {
    4: "Sionna's Path Solver launches rays isotropically from the transmitter, that is, "
       "evenly in every direction. When a ray hits a surface, the solver draws at random "
       "what happens there, a mirror reflection or a diffuse scatter, and the material's "
       "scattering coefficient sets the odds. The ray carries that material's reflection "
       "coefficient, and the result comes from the few paths that actually come back.\n\n"
       "Our kernel does something different. We shoot a grid of parallel rays at the target, "
       "take the first point each ray hits, and sum the physical optics scattering integral "
       "over those points. For parts we treat as transmissive, the ray passes through, "
       "and we add the echo from the metal inside, scaled by how much gets through.",

    5: "The centre frequency is 3.5 GHz, a standard 5G band. For the waveform I used a simple "
       "continuous wave. This is an early experiment, and CW is the easiest thing to work "
       "with. What I want to see first is the micro-Doppler itself, so I did not need a more "
       "complicated waveform yet.\n\n"
       "For the STFT: window 70 samples, about 3.6 milliseconds; hop 2 samples; "
       "FFT size 560.\n\n"
       "The bottom figure is made in two steps. First I add up the energy from 430 to "
       "1229 Hz, on both sides of zero. I keep away from zero on purpose, because the strong "
       "band there is mostly the body. What I keep is where the outer part of the blades "
       "lands. It is much weaker, but it belongs to the blades.\n\n"
       "Then I ask how that energy rises and falls over time, and take its spectrum.\n\n"
       "The peak sits at about 127 Hz. That is the blade flash rate: two blades passing per "
       "revolution, so twice the rotation rate. Its harmonics show up clearly too.\n\n"
       "[if asked] The band near zero moves as well, and it is much stronger. That is the "
       "body and the blades adding together, so it does carry the blade timing, but it is "
       "mixed with the body. I stay away from it for that reason.",

    6: "The previous slide was 3 metres only. Here I ran the same thing at 3, 15 and "
       "40 metres.\n\n"
       "Our kernel gives nearly the same map at every range. The Path Solver does not. "
       "One of the 40 metre runs looks different, and that turns out to be the random seed, "
       "not the distance. I will come back to it.",
}


def media(path: str) -> dict:
    z = zipfile.ZipFile(path)
    return {n: hashlib.md5(z.read(n)).hexdigest()
            for n in z.namelist() if n.startswith("ppt/media/")}


def main() -> None:
    if not os.path.exists(SRC):
        raise SystemExit(f"❌ {SRC} 가 없다")
    shutil.copy2(SRC, DST)
    prs = Presentation(DST)
    print(f"═══ v20({len(prs.slides)} 장) → v21 · 대본 {len(NOTES)} 개만 교체\n")

    for i, slide in enumerate(prs.slides, start=1):
        new = NOTES.get(i)
        if not new:
            continue
        old = (slide.notes_slide.notes_text_frame.text or "").strip()
        slide.notes_slide.notes_text_frame.text = new
        print(f"  [{i:>2}] {len(old):>4d} → {len(new):>4d} 자")
    prs.save(DST)

    # ── 검산 — 대본 셋 말고는 v20 과 같아야 한다 ────────────────────────────
    a, b = Presentation(SRC), Presentation(DST)
    same_body = all(
        [sh.text_frame.text for sh in sa.shapes if sh.has_text_frame]
        == [sh.text_frame.text for sh in sb.shapes if sh.has_text_frame]
        for sa, sb in zip(a.slides, b.slides))
    same_media = media(SRC) == media(DST)
    same_count = len(a.slides) == len(b.slides)
    notes_ok = all(b.slides[i - 1].notes_slide.notes_text_frame.text.strip() == v.strip()
                   for i, v in NOTES.items())
    untouched = [i for i in range(1, len(b.slides) + 1) if i not in NOTES]
    notes_kept = all(
        a.slides[i - 1].notes_slide.notes_text_frame.text.strip()
        == b.slides[i - 1].notes_slide.notes_text_frame.text.strip()
        for i in untouched if a.slides[i - 1].has_notes_slide)

    print(f"\n  장수 동일 {'✅' if same_count else '❌'}"
          f" · 본문 불변 {'✅' if same_body else '❌'}"
          f" · 그림 md5 전부 동일 {'✅' if same_media else '❌'}"
          f" · 새 대본 되읽기 {'✅' if notes_ok else '❌'}"
          f" · 나머지 대본 불변 {'✅' if notes_kept else '❌'}")
    print(f"\n✅ {DST}  ({os.path.getsize(DST):,} B)")
    if not (same_count and same_body and same_media and notes_ok and notes_kept):
        raise SystemExit("❌ 검산 실패 — 이 v21 을 쓰지 마라")


if __name__ == "__main__":
    main()
