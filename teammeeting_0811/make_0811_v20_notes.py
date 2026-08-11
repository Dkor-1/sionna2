# -*- coding: utf-8 -*-
"""
make_0811_v20_notes.py — **v19 의 슬라이드는 그대로 두고 대본만 갈아끼운다.**

사용자(2026-08-11)
> "전체적으로 대본 작성해서 v19 대본 수정해서 줄 수 있을까?"
> "간결성 유지하고 너무 길게 말하지 않으며 현학적인 표현을 과하지 않게"

⛔ v19 는 **안 건드린다**. 복사본 v20 을 새로 만든다.
⛔ 그림·본문·레이아웃은 한 글자도 안 바꾼다. `notes_slide` 만 교체한다.

■ 고친 것 다섯 (전부 **사실 오류**다. 다듬기는 안 넣었다)
 1. [3] "hard-coded **randomized** values" — 난수가 아니다. 원장 rpm_per_rotor 는
    3800 × [+0.220 %, −0.220 %, −0.121 %, +0.121 %] 로 손으로 정한 결정론적 패턴이다.
 2. [4] "under a **plane-wave** illumination assumption" — 덱의 거리별 그림은 구면파다
    (outputs/deck0811_range_figs.json: ours_illumination = "spherical wave, computed per range").
    파동 가정 문구를 빼고 «평행 광선 격자» 라는 사실만 남겼다.
 3. [5] "127 Hz ≈ **propeller rotation speed**" — 2 배 틀렸다. 126.67 Hz 는 **블레이드 플래시율**
    (날개 2 장 × 3,800 rpm ÷ 60)이고 회전수는 63.3 Hz 다.
 4. [5] "밴드 430–1229 Hz 인데 피크가 127 Hz" — 그대로 들으면 자기모순이다. 대역 전력을
    **시간축으로 다시 FFT** 한 것이라는 두 단계를 명시했다.
 5. [6] "patterns remained highly consistent" — 우리 커널에서만 참이다. PathSolver 40 m 는
    시드 하나가 2 배선에 앉는다([8] 에서 본인이 말할 내용과 충돌).

    /home/yunjung/.venvs/py312/bin/python teammeeting_0811/make_0811_v20_notes.py
"""
from __future__ import annotations

import os
import shutil

from pptx import Presentation

TM = "/home/yunjung/workspace/team_meeting"
SRC = f"{TM}/teammeeting_0811_v19.pptx"
DST = f"{TM}/teammeeting_0811_v20.pptx"

# 슬라이드 번호(1부터) → 대본. 빈 문자열이면 건드리지 않는다.
NOTES = {
    2: (
        "Today I show micro-Doppler maps of a hovering drone, computed two ways: "
        "with Sionna's Path Solver, and with our own SBR plus physical optics kernel."
    ),
    3: (
        "The target is a DJI Matrice 4E. Each rotor spins at a fixed speed near 3,800 RPM, "
        "which I got from a published thrust measurement. The four rotors differ slightly "
        "from each other, and that spread comes from Ji-hyuk's prior work. "
        "These are hand-set constant values, not measured from a real flight. "
        "Better rotor modelling is future work."
    ),
    4: (
        "Sionna's Path Solver spreads rays evenly over the full sphere from the transmitter. "
        "When a ray hits a surface it picks up that material's reflection coefficient, "
        "and the result comes from the few paths that actually come back.\n\n"
        "Our kernel does something different. We shoot a grid of parallel rays at the target, "
        "take the first point each ray hits, and sum the physical optics scattering integral "
        "over those points. For parts we treat as transmissive, the ray passes through, "
        "and we add the echo from the metal inside, scaled by how much gets through."
    ),
    5: (
        "The centre frequency is 3.5 GHz, a standard 5G band. I kept the waveform simple, "
        "a continuous wave. For the STFT: window 70 samples, about 3.6 milliseconds; "
        "hop 2 samples; FFT size 560.\n\n"
        "The bottom figure is made in two steps. First I take the energy in the 430 to 1229 Hz "
        "band, which is above the body and belongs to the blade tips. Then I ask how that "
        "energy rises and falls over time, and take its spectrum.\n\n"
        "The peak sits at about 127 Hz. That is the blade flash rate: two blades passing per "
        "revolution, so twice the rotation rate. Its harmonics show up clearly too."
    ),
    6: (
        "The previous slide was 3 metres only. Here I ran the same thing at 3, 15 and 40 metres.\n\n"
        "Our kernel gives nearly the same map at every range. The Path Solver does not. "
        "At 40 metres its map changes, and I will come back to why."
    ),
    7: (
        "This is the band energy for each range. One thing stands out in the harmonics. "
        "With our kernel the peaks fall off steadily as you go up. With the Path Solver, "
        "the second harmonic often comes out almost as strong as the fundamental, "
        "and sometimes stronger."
    ),
    8: (
        "At 40 metres I ran the Path Solver again with only the random seed changed. "
        "The map depends on the seed. One seed puts the strongest line on the second harmonic "
        "instead of the fundamental.\n\n"
        "Our kernel has no randomness. Same distance, same pose, same answer every time."
    ),
    9: (
        "For my future work, I want to make the simulation more realistic. "
        "I have three directions.\n\n"
        "First, the rotor speeds. Right now each rotor turns at a hard-coded value, and the "
        "spread between them comes from a flight simulator, not from a real drone. I want to "
        "use real flight log data instead, and also introduce temporal variation, the way a "
        "real drone adjusts its rotors while it stays in place.\n\n"
        "Second, noise. There is no receiver noise in this simulation. That is why nothing "
        "became harder when the drone moved further away. The echo does get weaker with "
        "distance, but each map is scaled to its own brightest point, so we cannot see that. "
        "And with no noise, nothing is ever buried. In the real world, noise is what decides "
        "how far we can detect, so I want to add it.\n\n"
        "Third, the geometry and the environment. Everything today uses one station that sends "
        "and receives at the same place, in empty space. Our real experiment is a passive "
        "bistatic setup, where the signal comes from someone else's base station, and it "
        "happens in a room with walls and a floor. I want to move to that geometry and add "
        "clutter modelling.\n\n"
        "All three will make the results worse, and that is the point. "
        "They bring the simulation closer to a real measurement."
    ),
}


def main() -> None:
    if not os.path.exists(SRC):
        raise SystemExit(f"❌ {SRC} 가 없다")
    shutil.copy2(SRC, DST)

    prs = Presentation(DST)
    n_slides = len(prs.slides)
    print(f"═══ v19 → v20 · 슬라이드 {n_slides} 장 (그림·본문 불변, 대본만 교체)\n")

    changed = 0
    for i, slide in enumerate(prs.slides, start=1):
        new = NOTES.get(i)
        if not new:
            print(f"  [{i}] (대본 없음 — 건드리지 않음)")
            continue
        old = (slide.notes_slide.notes_text_frame.text or "").strip()
        slide.notes_slide.notes_text_frame.text = new
        changed += 1
        print(f"  [{i}] {len(old):>4d} → {len(new):>4d} 자   {new.splitlines()[0][:62]}…")

    prs.save(DST)

    # ── 검산: 대본이 실제로 박혔나, 그리고 슬라이드 본문이 안 변했나 ──────────
    chk = Presentation(DST)
    ok = all((chk.slides[i - 1].notes_slide.notes_text_frame.text.strip() == v.strip())
             for i, v in NOTES.items())
    a, b = Presentation(SRC), chk
    same_body = all(
        [sh.text_frame.text for sh in sa.shapes if sh.has_text_frame]
        == [sh.text_frame.text for sh in sb.shapes if sh.has_text_frame]
        for sa, sb in zip(a.slides, b.slides))

    print(f"\n  대본 {changed} 장 교체 · 되읽기 일치 {'✅' if ok else '❌'}"
          f" · 본문 불변 {'✅' if same_body else '❌ 본문이 바뀌었다'}")
    print(f"\n✅ {DST}  ({os.path.getsize(DST):,} B)")
    if not (ok and same_body):
        raise SystemExit("❌ 검산 실패 — v20 을 쓰지 마라")


if __name__ == "__main__":
    main()
