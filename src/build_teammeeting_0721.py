"""팀미팅 0721 덱 빌더.

0714 덱(teammeeting_0714_v14.pptx)을 원본으로 열어 슬라이드만 갈아끼운다.
python-pptx 로는 23_ClassicWhite 마스터·레이아웃과 '/19' 페이지번호 placeholder 를
새로 만들 수 없기 때문에, 반드시 실제 파일에서 출발해야 한다.

슬라이드 내용은 content_0721.json 에서 읽는다 (검증된 수치만 담긴 파일).

사용:
    python src/build_teammeeting_0721.py <content.json> <out.pptx> [--src <0714.pptx>]
"""

from __future__ import annotations

import argparse
import copy
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ── 팔레트 (0714 에서 추출, 전부 하드코딩 srgbClr — 테마색 쓰지 말 것) ────────────
NAVY = RGBColor(0x26, 0x44, 0x8A)   # L1 불릿, chip 텍스트, Takeaway, DivRoad, 우측 카드
BLUE = RGBColor(0x2E, 0x74, 0xB5)   # Tag, 화살표, DivKicker, chip 테두리, 좌측 카드
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x59, 0x59, 0x59)   # 항상 이탤릭 — Cap, DivSub, DivMeta
RED = RGBColor(0xC0, 0x00, 0x00)    # Subtitle 전용. 프레이밍 한 줄, 덱 전체 4개 이하
FILL_L = RGBColor(0xEA, 0xF1, 0xFB)  # chip + 좌측 카드
FILL_R = RGBColor(0xE9, 0xED, 0xF5)  # 우측 카드 전용
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

F_BODY = "Avenir Next Medium"
F_HEAD = "AvenirNext-DemiBold"  # 공백 없음 — 별개 패밀리 문자열

A = qn("a:t").rsplit("}", 1)[0] + "}"  # '{...drawingml/2006/main}'


def _q(tag: str) -> str:
    return A + tag


# ── 저수준 헬퍼 ────────────────────────────────────────────────────────────────
def setfont(run, name=F_BODY, size=24, bold=False, italic=False, color=None):
    """python-pptx 의 font.name 은 latin 만 건드린다 — ea/cs 까지 직접 박는다."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        for old in rPr.findall(_q(tag)):
            rPr.remove(old)
        el = rPr.makeelement(_q(tag), {"typeface": name})
        rPr.append(el)


def box(slide, name, l, t, w, h, anchor=MSO_ANCHOR.MIDDLE):
    sh = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    sh.name = name
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(50800)
    tf.margin_top = tf.margin_bottom = Emu(30000)
    tf.vertical_anchor = anchor
    bodyPr = tf._txBody.find(_q("bodyPr"))
    for auto in ("normAutofit", "spAutoFit"):
        for old in bodyPr.findall(_q(auto)):
            bodyPr.remove(old)
    bodyPr.append(bodyPr.makeelement(_q("noAutofit"), {}))
    return sh


def put(sh, text, size, bold=False, italic=False, color=BLACK, name=F_BODY,
        align=PP_ALIGN.LEFT):
    tf = sh.text_frame
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    setfont(r, name, size, bold, italic, color)
    return sh


def _bullet_pPr(p, level):
    """L1: 40pt 굵은 네이비 + Arial '•'  /  L2: 34pt 검정 + Wingdings 'Ø'(➢)"""
    pPr = p._p.get_or_add_pPr()
    if level == 1:
        pPr.set("marL", "600000")
        pPr.set("indent", "-600000")
        spc, bu_font, bu_char = "1900", "Arial", "•"
    else:
        pPr.set("marL", "1260000")
        pPr.set("indent", "-660000")
        spc, bu_font, bu_char = "900", "Wingdings", "Ø"
    pPr.set("algn", "l")
    sb = pPr.makeelement(_q("spcBef"), {})
    pts = pPr.makeelement(_q("spcPts"), {"val": spc})
    sb.append(pts)
    pPr.append(sb)
    pPr.append(pPr.makeelement(_q("buFont"), {"typeface": bu_font}))
    pPr.append(pPr.makeelement(_q("buChar"), {"char": bu_char}))


def add_bullets(slide, items, l1_size=40, l2_size=34,
                l=720000, t=2820000, w=22944000, h=8400000):
    """items: ('L1'|'L2'|'', text). '' 는 그룹 구분용 빈 문단."""
    sh = box(slide, "Bullets", l, t, w, h, anchor=MSO_ANCHOR.TOP)
    tf = sh.text_frame
    first = True
    for level, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if not level:
            continue  # 빈 문단 = 그룹 구분자
        _bullet_pPr(p, 1 if level == "L1" else 2)
        r = p.add_run()
        r.text = text
        if level == "L1":
            setfont(r, F_BODY, l1_size, bold=True, color=NAVY)
        else:
            setfont(r, F_BODY, l2_size, bold=False, color=BLACK)
    return sh


def fill_line(sh, fill=None, line=None, width_emu=12700):
    if fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Emu(width_emu)
    else:
        sh.line.fill.background()


def add_pagenum(slide, total):
    """'빈 페이지' 레이아웃의 sldNum placeholder 를 복사해 온다 (add_slide 가 안 옮김)."""
    for ph in slide.placeholders:
        if ph.element.find(".//" + qn("p:ph")) is not None:
            f = ph.element.find(".//" + qn("p:ph"))
            if f.get("type") == "sldNum":
                return  # 이미 있음
    src = None
    for sp in slide.slide_layout.shapes:
        f = sp.element.find(".//" + qn("p:ph"))
        if f is not None and f.get("type") == "sldNum":
            src = sp
            break
    if src is None:
        return
    el = copy.deepcopy(src.element)
    slide.shapes._spTree.append(el)
    # '/19' 리터럴 런을 '/<total>' 로 교체
    for t in el.iter(_q("t")):
        if t.text and t.text.startswith("/"):
            t.text = "/%d" % total


# ── 아키타입 ──────────────────────────────────────────────────────────────────
def s_header(slide, tag, title, subtitle=""):
    if tag:
        put(box(slide, "Tag", 720000, 450000, 21000000, 600000),
            tag, 26, bold=True, color=BLUE)
    put(box(slide, "Title", 700000, 1040000, 22980000, 940000),
        title, 52, bold=True, color=BLACK, name=F_HEAD)
    if subtitle:
        put(box(slide, "Subtitle", 760000, 2040000, 22600000, 700000),
            subtitle, 33, bold=True, italic=True, color=RED)


def s_cap(slide, text):
    put(box(slide, "Cap", 760000, 10697676, 22944000, 560000),
        text, 20, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


def s_takeaway(slide, text):
    put(box(slide, "Takeaway", 720000, 11150000, 22944000, 850000),
        text, 36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def s_divider(slide, kicker, title, sub="", road="", meta=""):
    W = 24384000
    if kicker:
        put(box(slide, "DivKicker", 0, 4050000, W, 760000),
            kicker.upper(), 36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    put(box(slide, "DivTitle", 0, 4850000, W, 1550000),
        title, 72, bold=True, color=BLACK, name=F_HEAD, align=PP_ALIGN.CENTER)
    if sub:
        put(box(slide, "DivSub", 0, 6520000, W, 820000),
            sub, 40, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    if road:
        put(box(slide, "DivRoad", 0, 7780000, W, 820000),
            road, 36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    if meta:
        put(box(slide, "DivMeta", 0, 8950000, W, 600000),
            meta, 26, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


def s_chipgrid(slide, rows):
    """rows: ['라벨 | chip1 | chip2 | chip3', ...]  chip 은 'main;;sub' 로 2줄 가능."""
    XS = [5900000, 12082000, 17482000]
    CW, CH = 5400000, 1700000
    for i, row in enumerate(rows[:3]):
        parts = [p.strip() for p in row.split("|")]
        label, chips = parts[0], parts[1:]
        y = 3150000 + i * 2350000
        put(box(slide, "L", 720000, y, 4900000, CH), label, 25, bold=True, color=NAVY)
        for j, chip in enumerate(chips[:3]):
            if not chip:
                continue
            x = XS[j]
            if j > 0:
                put(box(slide, "A", x - 782000, y, 782000, CH),
                    "→", 36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
            sh = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(CW), Emu(CH))
            sh.name = "B"
            fill_line(sh, FILL_L, BLUE, 15875)
            tf = sh.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            main, _, sub = chip.partition(";;")
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = main
            setfont(r, F_BODY, 24, bold=True, color=NAVY)
            if sub:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                r2 = p2.add_run()
                r2.text = sub
                setfont(r2, F_BODY, 17, italic=True, color=GRAY)


def s_cards(slide, cards):
    """cards: ['LEFT|header|subhead|line1;line2', 'RIGHT|...']"""
    geo = {"LEFT": (820000, BLUE, FILL_L), "RIGHT": (12660000, NAVY, FILL_R)}
    for spec in cards:
        side, header, subhead, lines = (spec.split("|", 3) + ["", "", ""])[:4]
        side = side.strip().upper()
        if side not in geo:
            continue
        x, accent, fill = geo[side]
        hdr = slide.shapes.add_textbox(Emu(x), Emu(2980000), Emu(10900000), Emu(1120000))
        hdr.name = "CardHdr"
        fill_line(hdr, accent, None)
        tfh = hdr.text_frame
        tfh.word_wrap = True
        tfh.vertical_anchor = MSO_ANCHOR.MIDDLE
        ph = tfh.paragraphs[0]
        ph.alignment = PP_ALIGN.CENTER
        rh = ph.add_run()
        rh.text = header
        setfont(rh, F_HEAD, 35, bold=True, color=WHITE)

        bod = slide.shapes.add_textbox(Emu(x), Emu(4100000), Emu(10900000), Emu(6450000))
        bod.name = "CardBody"
        fill_line(bod, fill, accent, 12700)
        tfb = bod.text_frame
        tfb.word_wrap = True
        tfb.vertical_anchor = MSO_ANCHOR.TOP
        tfb.margin_left = tfb.margin_right = Emu(200000)
        tfb.margin_top = Emu(150000)
        first = True
        if subhead:
            p = tfb.paragraphs[0]
            r = p.add_run()
            r.text = subhead
            setfont(r, F_BODY, 35, bold=True, italic=True, color=accent)
            first = False
        for line in [x for x in lines.split(";") if x.strip()]:
            p = tfb.paragraphs[0] if first else tfb.add_paragraph()
            first = False
            r = p.add_run()
            r.text = line.strip()
            setfont(r, F_BODY, 32, color=BLACK)


def s_figure(slide, path, top=2800000, bottom=10500000):
    from PIL import Image
    if not os.path.exists(path):
        print(f"  !! 그림 없음: {path}")
        return
    with Image.open(path) as im:
        iw, ih = im.size
    band_h = bottom - top
    max_w = 22944000
    h = band_h
    w = int(h * iw / ih)
    if w > max_w:
        w = max_w
        h = int(w * ih / iw)
    left = int((24384000 - w) / 2)
    slide.shapes.add_picture(path, Emu(left), Emu(top + (band_h - h) // 2),
                             width=Emu(w), height=Emu(h))


# ── 빌드 ──────────────────────────────────────────────────────────────────────
def clear_slides(prs, keep=0):
    """keep 개만 남기고 전부 삭제 (0 이면 전부)."""
    xml_slides = prs.slides._sldIdLst
    ids = list(xml_slides)
    for sldId in ids[keep:]:
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)


def find_layout(prs, name):
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if lay.name == name:
                return lay
    return None


def build(content_path, out_path, src_path):
    content = json.load(open(content_path))
    slides = content["slides"]
    total = len(slides)

    prs = Presentation(src_path)
    assert prs.slide_width == 24384000 and prs.slide_height == 13716000, "슬라이드 크기 불일치"

    # 1번(타이틀)만 남기고 삭제 → 로고·Arial 랩블록을 그대로 물려받는다
    clear_slides(prs, keep=1)
    blank = find_layout(prs, "빈 페이지")
    assert blank is not None, "'빈 페이지' 레이아웃을 못 찾음"

    # 타이틀 슬라이드 날짜 갱신
    t0 = prs.slides[0]
    for sh in t0.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if "2026" in r.text and "07" in r.text:
                    r.text = content.get("date", "2026. 07. 21")

    for spec in slides[1:]:
        s = prs.slides.add_slide(blank)
        arch = spec["archetype"]
        tag = spec.get("tag", "")
        title = spec.get("title", "")
        sub = spec.get("subtitle", "")

        if arch == "divider":
            s_divider(s, spec.get("kicker", ""), title, spec.get("sub", ""),
                      spec.get("road", ""), spec.get("meta", ""))
        elif arch == "thanks":
            put(box(s, "DivTitle", 0, 6083000, 24384000, 1550000),
                title or "Thank you!", 100, bold=True, color=BLACK,
                name=F_HEAD, align=PP_ALIGN.CENTER)
        else:
            s_header(s, tag, title, sub)
            if arch in ("bullets", "recap", "contents"):
                items = []
                for b in spec.get("bullets", []):
                    if b.strip() == "":
                        items.append(("", ""))
                    elif b.startswith("L1:"):
                        items.append(("L1", b[3:].strip()))
                    elif b.startswith("L2:"):
                        items.append(("L2", b[3:].strip()))
                    else:
                        items.append(("L2", b.strip()))
                sz = (48, 40) if arch == "contents" else (40, 34)
                add_bullets(s, items, sz[0], sz[1])
            elif arch in ("chipgrid", "status"):
                s_chipgrid(s, spec.get("chips", []))
            elif arch == "cards":
                s_cards(s, spec.get("cards", []))
            elif arch == "figure":
                fig = spec.get("figure", "")
                if fig:
                    s_figure(s, fig)
                if spec.get("cap"):
                    s_cap(s, spec["cap"])
            if spec.get("takeaway"):
                if arch == "figure":
                    print(f"  !! 규약 위반: figure 슬라이드({spec['n']})에 Takeaway — 무시함")
                else:
                    s_takeaway(s, spec["takeaway"])

        add_pagenum(s, total)

        if spec.get("notes"):
            s.notes_slide.notes_text_frame.text = spec["notes"]

    prs.save(out_path)
    print(f"저장: {out_path}  ({total} 슬라이드)")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("content")
    ap.add_argument("out")
    ap.add_argument("--src", default="/workspace/team_meeting/teammeeting_0714_v14.pptx")
    a = ap.parse_args()
    build(a.content, a.out, a.src)


if __name__ == "__main__":
    main()
